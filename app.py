import os
import django
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from pydantic import BaseModel, HttpUrl
from typing import Optional, List
from langserve import add_routes
from langchain_core.runnables import RunnablePassthrough
from config import llm, embeddings, streaming_llm
from langchain.text_splitter import RecursiveCharacterTextSplitter
from fastapi.responses import StreamingResponse
from typing import List
import asyncio
from logger import log_manager

# 设置Django环境（仅用于ORM）
os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'settings')
django.setup()

# 创建FastAPI应用
app = FastAPI(
    title="AI编辑器助手API",
    description="基于LangChain的AI编辑器助手API服务",
    version="1.0.0"
)

# 在Django初始化完成后导入和创建AI编辑器助手实例
from ai_editor import AIEditorAssistant
assistant = AIEditorAssistant(llm, streaming_llm, embeddings)

# 请求模型
class TextGenerationRequest(BaseModel):
    user_text: Optional[str] = None
    prompt: str
    file: Optional[UploadFile] = None
    url: Optional[HttpUrl] = None

# 创建一个Runnable对象来包装generate_text方法
generate_runnable = (
    RunnablePassthrough()
    | {
        "user_text": lambda x: x.get("user_text", ""),
        "prompt": lambda x: x["prompt"],
        "result": lambda x: assistant.generate_text(x.get("user_text", ""), x["prompt"])
    }
)

class UserEditRequest(BaseModel):
    original_text: str
    edited_text: str
    text_type: Optional[str] = None

# 添加API路由
@app.post("/generate", response_model=str)
async def generate_text(request: TextGenerationRequest):
    """生成文本内容"""
    try:
        # 确保user_text不为None，如果为None则使用空字符串
        return assistant.generate_text(request.user_text or "", request.prompt)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/record-edit")
async def record_edit(request: UserEditRequest):
    """记录用户编辑"""
    try:
        assistant.record_user_edit(
            request.original_text,
            request.edited_text,
            request.text_type
        )
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/preferences/{text_type}")
async def get_preferences(text_type: str = "general"):
    """获取用户偏好"""
    try:
        return {"preferences": assistant.preference_manager.get_preferences(text_type)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/test-log")
@log_manager.auto_log_request
async def test_log_endpoint():
    """测试日志装饰器功能的端点"""
    try:
        # 模拟一些操作
        result = {"message": "这是一个测试消息", "status": "success"}
        # 随机抛出异常来测试错误日志
        import random
        if random.random() < 0.5:  # 50%的概率抛出异常
            raise ValueError("这是一个测试异常")
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-with-context")
@log_manager.auto_log_request
async def generate_with_context(
    prompt: str = Form(...),
    user_text: Optional[str] = Form(None),
    files: Optional[List[UploadFile]] = File(None),
    urls: Optional[List[HttpUrl]] = Form(None)
):
    """生成带有文件或URL上下文的文本内容（流式输出）"""
    try:
        context = ""
        
        # 处理多个文件
        if files:
            for file in files:
                content = await file.read()
                file_extension = file.filename.split('.')[-1].lower()
                file_context = ""
                if file_extension in ['txt', 'md']:
                    file_context = content.decode('utf-8')
                elif file_extension == 'pdf':
                    from PyPDF2 import PdfReader
                    import io
                    reader = PdfReader(io.BytesIO(content))
                    file_context = '\n'.join(page.extract_text() for page in reader.pages)
                elif file_extension in ['docx', 'doc']:
                    from docx import Document
                    import io
                    doc = Document(io.BytesIO(content))
                    
                    # 提取段落文本
                    paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
                    
                    # 提取表格内容
                    tables = []
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                tables.append(row_text)
                    
                    # 提取列表内容
                    lists = [f"• {para.text}" for para in doc.paragraphs if para._p.pPr and para._p.pPr.numPr]
                    
                    file_context = '\n'.join(paragraphs + tables + lists)
                    
                elif file_extension in ['ppt', 'pptx']:
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(content))
                    slides_text = []
                    for slide in prs.slides:
                        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                        slides_text.extend(slide_text)
                    file_context = '\n'.join(slides_text)
                else:
                    continue  # 跳过不支持的文件类型
                
                context += f"\n=== 文件：{file.filename} ===\n{file_context}\n"
            
        # 并行处理多个URL
        if urls:
            import aiohttp
            from bs4 import BeautifulSoup
            
            async def fetch_url(url):
                async with aiohttp.ClientSession() as session:
                    try:
                        async with session.get(str(url)) as response:
                            if response.status == 200:
                                html = await response.text()
                                soup = BeautifulSoup(html, 'html.parser')
                                return f"\n=== URL：{url} ===\n{soup.get_text()}\n"
                    except Exception:
                        return ""
                return ""
            
            url_contexts = await asyncio.gather(*[fetch_url(url) for url in urls])
            context += ''.join(url_contexts)

        if context:
            # 使用文本分割器将文档分成合适大小的块
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
                separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?"]
            )
            text_chunks = text_splitter.split_text(context)
            # 创建临时向量存储
            from langchain_community.vectorstores import FAISS
            temp_vectorstore = FAISS.from_texts(text_chunks, embeddings)
        else:
            temp_vectorstore = None
        
        # 使用异步生成器进行流式输出
        async def generate_stream():
            try:
                async for chunk in assistant.generate_text_with_temp_context(
                    user_text=user_text,
                    prompt=prompt,
                    temp_vectorstore=temp_vectorstore,
                    stream=True
                ):
                    # 将每个文本块格式化为SSE消息格式
                    yield f"data: {str(chunk)}\n\n"
            except Exception as e:
                yield f"data: {str(e)}\n\n"  # 错误处理

        return StreamingResponse(
            generate_stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Content-Encoding": "none"
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# 创建另一个Runnable对象来包装generate_with_context方法
generate_with_context_runnable = (
    RunnablePassthrough()
    | {
        "user_text": lambda x: x.get("user_text", ""),
        "prompt": lambda x: x["prompt"],
        "file": lambda x: x.get("file", None),
        "url": lambda x: x.get("url", None),
        "result": lambda x: assistant.generate_text_with_temp_context(
            user_text=x.get("user_text", ""),
            prompt=x["prompt"],
            temp_vectorstore=None  # 需要根据file/url内容创建
        )
    }
)

# 添加LangServe路由
add_routes(
    app,
    generate_runnable,
    path="/langserve/generate",
    input_type=TextGenerationRequest,
    config_keys=["user_text", "prompt"]
)

add_routes(
    app,
    generate_with_context_runnable,
    path="/langserve/generate-with-context",
    input_type=TextGenerationRequest,
    config_keys=["user_text", "prompt", "file", "url"]
)

if __name__ == "__main__":
    import uvicorn
    # 确保日志管理器已初始化
    uvicorn.run("app:app", host="0.0.0.0", port=8000, reload=True)
    



    

